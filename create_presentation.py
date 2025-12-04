"""
Script to create a professional PowerPoint presentation
for the Sentiment Analysis project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(12)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    
    # Add image if it exists
    if Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(2), width=Inches(7))
        
        # Add caption if provided
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].font.size = Pt(14)
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_table_slide(prs, title, data_dict):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    
    # Create table
    rows = len(data_dict['data']) + 1  # +1 for header
    cols = len(data_dict['headers'])
    
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
    
    # Set column headers
    for col_idx, header in enumerate(data_dict['headers']):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 121)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Fill data
    for row_idx, row_data in enumerate(data_dict['data'], start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    return slide

# ===== SLIDE 1: Title =====
add_title_slide(
    prs,
    "Analyse Automatique des Sentiments",
    "Projet NLP - Dataset Sentiment140"
)

# ===== SLIDE 2: Introduction =====
add_content_slide(
    prs,
    "📋 Introduction",
    [
        "🎯 Objectif: Classifier automatiquement les sentiments dans les tweets",
        "📊 Dataset: Sentiment140 (1.6 million de tweets)",
        "⚖️ Classes: Positif (😊) vs Négatif (😞)",
        "🛠️ Technologie: Python, Scikit-learn, NLTK, Pandas",
        "📅 Période: Semaines 1-3 du planning projet"
    ]
)

# ===== SLIDE 3: Dataset Overview =====
add_content_slide(
    prs,
    "📊 Dataset: Sentiment140",
    [
        "✓ Source: Stanford University (Twitter data)",
        "✓ Taille: 1,600,000 tweets",
        "✓ Distribution: 800K négatifs (0) + 800K positifs (4)",
        "✓ Balance parfaite: 50% / 50%",
        "✓ Colonnes: target, text, user, date, ids",
        "✓ Fichier: 73.2 MB après nettoyage",
        "✓ Échantillon utilisé: 100,000 tweets (50K+50K)"
    ]
)

# ===== SLIDE 4: Semaine 1 - Exploration =====
add_content_slide(
    prs,
    "🔍 Semaine 1: Exploration Initiale",
    [
        "✓ Téléchargement du dataset depuis Kaggle",
        "✓ Chargement avec encoding latin-1",
        "✓ Analyse de la structure des données",
        "✓ Vérification des valeurs manquantes: 0",
        "✓ Statistiques textuelles:",
        "   • Longueur moyenne: 74 caractères",
        "   • Nombre moyen de mots: 13 mots/tweet",
        "   • Distribution équilibrée confirmée"
    ]
)

# ===== SLIDE 5: Semaine 2 - Prétraitement (1) =====
add_content_slide(
    prs,
    "🧹 Semaine 2: Prétraitement des Données",
    [
        "✓ Nettoyage du texte:",
        "   • Conversion en minuscules",
        "   • Suppression des URLs (http://, https://)",
        "   • Suppression des mentions (@user)",
        "   • Suppression des hashtags (#topic)",
        "   • Suppression des caractères spéciaux et chiffres",
        "   • Normalisation des espaces",
        "✓ Résultat: texte propre prêt pour la modélisation"
    ]
)

# ===== SLIDE 6: Visualisations =====
add_content_slide(
    prs,
    "📈 Visualisations - Analyses Effectuées",
    [
        "✓ Distribution des sentiments (bar chart + pie chart)",
        "✓ Histogrammes des longueurs de texte par sentiment",
        "✓ Distribution du nombre de mots par sentiment",
        "✓ Word Clouds pour négatif et positif",
        "✓ Top 20 mots les plus fréquents par sentiment",
        "✓ Toutes les visualisations sauvegardées en haute résolution"
    ]
)

# ===== SLIDE 7: Word Clouds =====
# Try to add word cloud image if it exists
wordcloud_path = "visuals/wordclouds/sentiment_wordclouds.png"
add_image_slide(
    prs,
    "☁️ Word Clouds: Mots Fréquents",
    wordcloud_path,
    "Mots les plus fréquents dans les tweets négatifs (rouge) et positifs (vert)"
)

# ===== SLIDE 8: Semaine 3 - ML Models =====
add_content_slide(
    prs,
    "🤖 Semaine 3: Modèles ML Baseline",
    [
        "✓ Vectorisation: TF-IDF (10,000 features)",
        "   • Term Frequency - Inverse Document Frequency",
        "   • Unigrams + Bigrams (1-2 mots)",
        "   • min_df=5, max_df=0.7",
        "✓ Split: Train (80%) / Test (20%)",
        "   • Train: 79,788 tweets",
        "   • Test: 19,947 tweets",
        "✓ Stratification pour maintenir la balance"
    ]
)

# ===== SLIDE 9: 3 Modèles Entraînés =====
add_content_slide(
    prs,
    "🎯 Modèles ML Implémentés",
    [
        "1️⃣ Logistic Regression",
        "   • Modèle linéaire simple et rapide",
        "   • Bon baseline pour classification binaire",
        "",
        "2️⃣ Support Vector Machine (LinearSVC)",
        "   • Trouve le meilleur hyperplan de séparation",
        "   • Excellent pour classification de texte",
        "",
        "3️⃣ Random Forest",
        "   • Ensemble de 100 arbres de décision",
        "   • Capture les patterns non-linéaires"
    ]
)

# ===== SLIDE 10: Résultats Comparatifs =====
results_data = {
    'headers': ['Modèle', 'Accuracy', 'Precision', 'Recall', 'F1-Score'],
    'data': [
        ['Logistic Regression', '79.19%', '0.7922', '0.7915', '0.7918'],
        ['SVM (LinearSVC)', '79.15%', '0.7918', '0.7912', '0.7915'],
        ['Random Forest', '77.24%', '0.7728', '0.7721', '0.7724']
    ]
}

add_table_slide(
    prs,
    "📊 Résultats: Comparaison des Modèles",
    results_data
)

# ===== SLIDE 11: Métriques Expliquées =====
add_content_slide(
    prs,
    "📐 Métriques d'Évaluation",
    [
        "✓ Accuracy: Proportion de prédictions correctes",
        "   • (TP + TN) / Total",
        "",
        "✓ Precision: Parmi les prédictions positives, combien sont correctes",
        "   • TP / (TP + FP)",
        "",
        "✓ Recall: Parmi les vrais positifs, combien sont détectés",
        "   • TP / (TP + FN)",
        "",
        "✓ F1-Score: Moyenne harmonique de Precision et Recall",
        "   • 2 × (Precision × Recall) / (Precision + Recall)"
    ]
)

# ===== SLIDE 12: Matrices de Confusion =====
add_content_slide(
    prs,
    "🎯 Matrices de Confusion",
    [
        "Analyse des erreurs pour chaque modèle:",
        "",
        "✓ True Positives (TP): Positifs correctement prédits",
        "✓ True Negatives (TN): Négatifs correctement prédits",
        "✓ False Positives (FP): Négatifs prédits comme positifs",
        "✓ False Negatives (FN): Positifs prédits comme négatifs",
        "",
        "📊 3 matrices générées et sauvegardées",
        "📈 ~79% de précision globale"
    ]
)

# ===== SLIDE 13: Confusion Matrix Image =====
cm_path = "visuals/confusion_matrices/logistic_regression_cm.png"
add_image_slide(
    prs,
    "📊 Matrice de Confusion - Logistic Regression",
    cm_path,
    "Meilleur modèle avec 79.19% d'accuracy"
)

# ===== SLIDE 14: Exemple de Prédictions =====
add_content_slide(
    prs,
    "✅ Tests sur Exemples Réels",
    [
        "Tweet: \"i love this product it is amazing\"",
        "→ Prédiction: 😊 POSITIVE ✓",
        "",
        "Tweet: \"this is the worst experience ever\"",
        "→ Prédiction: 😞 NEGATIVE ✓",
        "",
        "Tweet: \"great quality highly recommend\"",
        "→ Prédiction: 😊 POSITIVE ✓",
        "",
        "Tweet: \"waste of money do not buy\"",
        "→ Prédiction: 😞 NEGATIVE ✓"
    ]
)


# ===== SLIDE 16: Points Clés =====
add_content_slide(
    prs,
    "🌟 Points Clés & Réussites",
    [
        "✅ Dataset équilibré: pas de biais de classe",
        "✅ Prétraitement robuste: texte propre et normalisé",
        "✅ TF-IDF efficace: 10K features pertinentes",
        "✅ 3 modèles ML implémentés et comparés",
        "✅ 79% accuracy: bon résultat pour un baseline",
        "✅ Tous les modèles sauvegardés pour réutilisation",
        "✅ Visualisations complètes et professionnelles",
        "✅ Code documenté et reproductible"
    ]
)

# ===== SLIDE 17: Prochaines Étapes =====
add_content_slide(
    prs,
    "🚀 Prochaines Étapes (Semaines 4-7)",
    [
        "📅 Semaine 4: Deep Learning",
        "   • RNN, LSTM, GRU pour capturer le contexte séquentiel",
        "",
        "📅 Semaine 5: Transfer Learning",
        "   • BERT fine-tuning pour améliorer les performances",
        "   • Hyperparameter tuning avec GridSearchCV",
        "",
        "📅 Semaine 6: Clustering",
        "   • Analyse non-supervisée (K-Means, DBSCAN)",
        "",
        "📅 Semaine 7: Finalisation",
        "   • Rapport final et présentation"
    ]
)

# ===== SLIDE 18: Ce qui Reste à Faire (Section 3.2) =====
add_content_slide(
    prs,
    "📝 Section 3.2 - Éléments à Compléter",
    [
        "Pour compléter la section 3.2 du cahier des charges:",
        "",
        "🔲 Bag of Words (BoW) - Autre méthode de vectorisation",
        "🔲 Word Embeddings (Word2Vec, GloVe)",
        "🔲 GridSearchCV / RandomizedSearchCV",
        "   • Optimisation automatique des hyperparamètres",
        "🔲 AUC-ROC comme métrique supplémentaire",
        "🔲 Analyse approfondie des erreurs",
        "   • Identifier les phrases difficiles à classer",
        "",
        "→ Prévu dans Notebook 02B"
    ]
)

# ===== SLIDE 19: Technologies & Outils =====
add_content_slide(
    prs,
    "🛠️ Technologies Utilisées",
    [
        "🐍 Python 3.11",
        "📓 Jupyter Notebook",
        "📊 Pandas & NumPy (manipulation de données)",
        "🤖 Scikit-learn (modèles ML)",
        "📈 Matplotlib & Seaborn (visualisations)",
        "💬 NLTK (traitement du langage naturel)",
        "☁️ WordCloud (nuages de mots)",
        "💾 Pickle (sauvegarde des modèles)",
        "📁 Git & GitHub (versioning)"
    ]
)

# ===== SLIDE 20: Conclusion =====
add_content_slide(
    prs,
    "🎓 Conclusion",
    [
        "✅ Semaines 1-3 complétées avec succès",
        "",
        "📊 Dataset exploré et bien compris",
        "🧹 Prétraitement efficace et reproductible",
        "🤖 Baseline ML solide: 79% accuracy",
        "📈 Visualisations claires et informatives",
        "",
        "🎯 Objectif atteint pour cette phase",
        "🚀 Prêt pour les modèles Deep Learning",
        "",
        "📖 Code disponible sur GitHub",
        "📧 Questions ?"
    ]
)

# ===== SLIDE 21: Merci =====
add_title_slide(
    prs,
    "Merci pour votre attention ! 🙏",
    "Questions ?"
)

# Save presentation
output_path = Path("reports/Presentation_Sentiment_Analysis.pptx")
output_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_path))

print(f"✅ Présentation créée avec succès !")
print(f"📁 Fichier: {output_path}")
print(f"📊 Nombre de slides: {len(prs.slides)}")
print(f"\n🎯 La présentation contient:")
print(f"   • Slide de titre")
print(f"   • Introduction et objectifs")
print(f"   • Description du dataset")
print(f"   • Travail des semaines 1-3")
print(f"   • Résultats comparatifs des modèles")
print(f"   • Visualisations et matrices de confusion")
print(f"   • Prochaines étapes")
print(f"   • Technologies utilisées")
print(f"   • Conclusion")
